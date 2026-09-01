"""
Academix AI — File Parser Utilities

Extracts text from uploaded files (PDF, DOCX, PPTX, TXT).
Used by the RAG ingestion pipeline to process uploaded course materials.
"""

import io
from typing import Optional


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF file."""
    from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    text_parts = []
    for page_num, page in enumerate(reader.pages, 1):
        page_text = page.extract_text()
        if page_text:
            text_parts.append(f"[Page {page_num}]\n{page_text}")
    return "\n\n".join(text_parts)


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from a DOCX file."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
    return "\n\n".join(text_parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from a PPTX file."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if slide_text:
            text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
    return "\n\n".join(text_parts)


def extract_text(file_bytes: bytes, file_name: str) -> str:
    """
    Extract text from a file based on its extension.
    
    Supported formats: PDF, DOCX, PPTX, TXT
    Returns the extracted text as a single string.
    """
    ext = file_name.lower().rsplit(".", 1)[-1] if "." in file_name else ""

    if ext == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == "docx":
        return extract_text_from_docx(file_bytes)
    elif ext == "pptx":
        return extract_text_from_pptx(file_bytes)
    elif ext in ("txt", "md", "text"):
        return file_bytes.decode("utf-8", errors="replace")
    else:
        raise ValueError(f"Unsupported file format: .{ext}. Supported: PDF, DOCX, PPTX, TXT")


def get_page_ref(text: str, chunk_start: int) -> Optional[str]:
    """
    Extract page/slide reference from chunk text.
    Looks for [Page X] or [Slide X] markers.
    """
    # Search backwards from chunk_start for the nearest page marker
    search_text = text[:chunk_start + 200]
    
    # Find last occurrence of page marker
    for marker in ["[Page ", "[Slide "]:
        idx = search_text.rfind(marker)
        if idx != -1:
            end_idx = search_text.find("]", idx)
            if end_idx != -1:
                return search_text[idx + 1:end_idx]
    return None
